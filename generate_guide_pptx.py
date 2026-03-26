#!/usr/bin/env python3
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

DK1       = RGBColor(0x26, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DK2       = RGBColor(0x11, 0x56, 0x7F)
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)
TEAL      = RGBColor(0x71, 0xD3, 0xDC)
ORANGE    = RGBColor(0xFF, 0x9F, 0x36)
BODY_GREY = RGBColor(0x5B, 0x5B, 0x5B)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)

TEMPLATE_SEARCH = [
    os.path.join(os.getcwd(), "templates", "snowflake_template.pptx"),
    os.path.expanduser("~/.cortex/skills/CoCo_pptx_Skill/snowflake_template.pptx"),
]
TEMPLATE = next((p for p in TEMPLATE_SEARCH if os.path.isfile(p)), None)
assert TEMPLATE, "snowflake_template.pptx not found"

prs = Presentation(TEMPLATE)

while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = (sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
           or sldId.get('r:id'))
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

def set_ph(slide, idx, text):
    ph = slide.placeholders[idx]
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def set_ph_lines(slide, idx, lines, font_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if not line.strip():
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if font_size:
            p.font.size = Pt(font_size)

def set_ph_sections(slide, idx, sections, heading_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    first = True
    for heading, body_lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        first = False
        run = p.add_run()
        run.text = heading
        run.font.bold = True
        run.font.color.rgb = DK2
        if heading_size:
            run.font.size = Pt(heading_size)
        for line in body_lines:
            bp = tf.add_paragraph()
            bp.level = 1
            bp.text = line
            if body_size:
                bp.font.size = Pt(body_size)

def add_shape_text(slide, shape_type, x, y, w, h, text, fill_color, text_color, font_size=10, bold=False):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    return shape

# === SLIDE 1: COVER ===
slide = prs.slides.add_slide(prs.slide_layouts[13])
set_ph(slide, 3, "WORKSHOP\nSNOWFLAKE")
set_ph(slide, 0, "Mediaset — Hands-on Lab")
set_ph(slide, 2, "Snowflake Solution Engineering  |  2026")

# === SLIDE 2: AGENDA ===
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "AGENDA")
set_ph(slide, 1, "3.5 ore di laboratorio pratico")

agenda_items = [
    ("40 min", "MODULO 1", "Setup e Basi\nSnowflake"),
    ("35 min", "MODULO 2", "Row & Column\nSecurity"),
    ("30 min", "MODULO 3", "Dynamic\nTables"),
    ("25 min", "MODULO 4", "Cortex AI\nSQL"),
    ("40 min", "MODULO 5", "Semantic Views\n& Analyst"),
    ("25 min", "MODULO 6", "Cortex Search\n& Intelligence"),
]
n = len(agenda_items)
box_w = (9.10 - 0.15 * (n - 1)) / n
x = 0.40

for i, (time, module, topic) in enumerate(agenda_items):
    t_box = slide.shapes.add_textbox(Inches(x), Inches(1.30), Inches(box_w), Inches(0.25))
    t_box.text_frame.word_wrap = True
    p = t_box.text_frame.paragraphs[0]
    p.text = time
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = BODY_GREY
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    bg = SF_BLUE if i % 2 == 0 else DK2
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.60, box_w, 0.50, module, bg, WHITE, 10, True)
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.20, box_w, 0.70, topic, LIGHT_BG, DK1, 9, False)
    x += box_w + 0.15

# === MODULE 1: SETUP ===
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "MODULO 1\nSETUP E BASI")

slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "SNOWFLAKE — CONCETTI FONDAMENTALI")
set_ph(slide, 2, "Piattaforma cloud-native con separazione storage/compute")
set_ph_sections(slide, 1, [
    ("Database e Schema", [
        "Organizzano logicamente i dati",
        "Schema: RAW, ANALYTICS, SECURITY",
    ]),
    ("Warehouse", [
        "Risorse di calcolo elastiche on-demand",
        "Auto-suspend e auto-resume",
        "Dimensioni: XSMALL → 6XLARGE",
    ]),
    ("Ruoli", [
        "Controllano accesso con least privilege",
        "SYSADMIN per oggetti, SECURITYADMIN per ruoli",
        "Custom roles riportano a SYSADMIN",
    ]),
], heading_size=12, body_size=10)

slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "GERARCHIA DEI RUOLI")
set_ph(slide, 1, "Best practice: tutti i ruoli custom riportano a SYSADMIN")

roles = [
    ("ACCOUNTADMIN", 0.40, 1.40, DK2),
    ("SYSADMIN", 2.20, 2.20, DK2),
    ("SECURITYADMIN", 5.50, 2.20, DK2),
    ("MEDIASET_ADMIN", 2.20, 3.20, SF_BLUE),
    ("MEDIASET_ANALYST", 1.20, 4.20, TEAL),
    ("MEDIASET_MARKETING", 3.20, 4.20, TEAL),
]
for name, x, y, color in roles:
    txt_color = WHITE if color in [DK2, SF_BLUE] else DK1
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 1.80, 0.50, name, color, txt_color, 9, True)

# === MODULE 2: SECURITY ===
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "MODULO 2\nSECURITY")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "ROW & COLUMN LEVEL SECURITY")
set_ph(slide, 3, "Protezione dati sensibili senza modificare le query")
set_ph_sections(slide, 1, [
    ("Masking Policy", [
        "Nasconde/offusca dati PII per colonna",
        "Basata sul ruolo dell'utente",
        "Email, telefono, importi finanziari",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("Row Access Policy", [
        "Filtra automaticamente le righe visibili",
        "Ogni utente vede solo i suoi dati",
        "Esempio: dati regionali Nord/Sud",
    ]),
], heading_size=11, body_size=10)

slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "ESEMPIO: EMAIL MASKING")
set_ph(slide, 1, "Comportamento diverso per ruolo")

examples = [
    ("ADMIN", "marco.rossi@email.it", DK2, WHITE),
    ("MARKETING", "ma***@email.it", SF_BLUE, WHITE),
    ("ANALYST", "***RISERVATO***", BODY_GREY, WHITE),
]
x = 0.40
for role, email, color, txt_color in examples:
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.50, 2.90, 0.40, role, color, txt_color, 11, True)
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.00, 2.90, 0.50, email, LIGHT_BG, DK1, 10, False)
    x += 3.10

# === MODULE 3: DYNAMIC TABLES ===
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "MODULO 3\nDYNAMIC TABLES")

slide = prs.slides.add_slide(prs.slide_layouts[5])
set_ph(slide, 0, "DYNAMIC TABLES — PIPELINE DECLARATIVE")
set_ph(slide, 2, "Definisci la trasformazione, Snowflake gestisce il refresh")
set_ph_sections(slide, 1, [
    ("Caratteristiche", [
        "Aggiornamento automatico incrementale",
        "TARGET_LAG: frequenza di refresh",
        "Pipeline a cascata (DT che dipendono da altre DT)",
    ]),
    ("Casi d'uso", [
        "Aggregazioni real-time (ascolti giornalieri)",
        "KPI pre-calcolati (top programmi settimanali)",
        "ETL senza orchestrazione esterna",
    ]),
], heading_size=12, body_size=10)

# === MODULE 4: CORTEX AI ===
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "MODULO 4\nCORTEX AI SQL")

slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "CORTEX AI FUNCTIONS")
set_ph(slide, 1, "LLM integrati direttamente nelle query SQL")

functions = [
    ("SENTIMENT", "Analizza sentiment\ndei testi", SF_BLUE),
    ("CLASSIFY", "Classifica in\ncategorie", DK2),
    ("SUMMARIZE", "Riassume testi\nlunghi", SF_BLUE),
    ("TRANSLATE", "Traduce tra\nlingue", DK2),
    ("COMPLETE", "Genera testo\ncon LLM", SF_BLUE),
]
x = 0.40
box_w = 1.70
for name, desc, color in functions:
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.50, box_w, 0.50, name, color, WHITE, 10, True)
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.10, box_w, 0.70, desc, LIGHT_BG, DK1, 9, False)
    x += box_w + 0.15

# === MODULE 5: SEMANTIC VIEWS ===
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "MODULO 5\nSEMANTIC VIEWS")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "SEMANTIC VIEWS & CORTEX ANALYST")
set_ph(slide, 3, "Metriche di business + query in linguaggio naturale")
set_ph_sections(slide, 1, [
    ("Semantic Views", [
        "Definizioni centralizzate di metriche",
        "Facts, Dimensions, Metrics",
        "Sinonimi per termini di business",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("Cortex Analyst", [
        "Query in linguaggio naturale",
        "Traduce domande in SQL corretto",
        "Usa le definizioni semantic view",
    ]),
], heading_size=11, body_size=10)

slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "ESEMPI DI DOMANDE")
set_ph(slide, 1, "Cortex Analyst risponde in linguaggio naturale")

questions = [
    "Qual è il programma con lo share più alto?",
    "Mostrami i top 5 programmi in prime time",
    "Confronta gli ascolti di Canale 5 e Italia 1",
    "Quanto budget pubblicitario dal settore alimentare?",
]
y = 1.50
for q in questions:
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.40, y, 9.10, 0.45, f"\"  {q}  \"", LIGHT_BG, DK1, 11, False)
    y += 0.55

# === MODULE 6: CORTEX SEARCH ===
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "MODULO 6\nCORTEX SEARCH")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_ph(slide, 0, "CORTEX SEARCH & SNOWFLAKE INTELLIGENCE")
set_ph(slide, 3, "Ricerca semantica e assistente AI integrato")
set_ph_sections(slide, 1, [
    ("Cortex Search", [
        "Ricerca per significato, non keyword",
        "Trova contenuti semanticamente simili",
        "Ideale per RAG e recommendation",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("Snowflake Intelligence", [
        "Assistente AI integrato in Snowsight",
        "Combina Analyst + Search",
        "Risponde su tutti i tuoi dati",
    ]),
], heading_size=11, body_size=10)

# === SUMMARY ===
slide = prs.slides.add_slide(prs.slide_layouts[7])
set_ph(slide, 0, "KEY TAKEAWAYS")
set_ph(slide, 4, "Cosa abbiamo imparato oggi")
set_ph_sections(slide, 1, [
    ("Fondamenti", [
        "Database, Schema, Warehouse",
        "Ruoli e permessi",
        "Gerarchia SYSADMIN",
    ]),
], heading_size=10, body_size=9)
set_ph_sections(slide, 2, [
    ("Security & Pipeline", [
        "Masking e Row Access Policy",
        "Dynamic Tables",
        "ETL declarativo",
    ]),
], heading_size=10, body_size=9)
set_ph_sections(slide, 3, [
    ("AI & Analytics", [
        "Cortex AI Functions",
        "Semantic Views",
        "Cortex Search",
    ]),
], heading_size=10, body_size=9)

# === THANK YOU ===
slide = prs.slides.add_slide(prs.slide_layouts[28])
set_ph(slide, 1, "GRAZIE!")

# === SAVE ===
output_path = os.path.expanduser("~/Downloads/Mediaset_Workshop_Guide.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
